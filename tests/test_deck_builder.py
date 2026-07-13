import os

import duckdb
import pytest
from pptx import Presentation

from reporting.deck_builder import build_deck

TEST_DB = os.path.join(os.path.dirname(__file__), "..", "data", "ci_test.duckdb")
DB_PATH = os.environ.get("DBT_DB_PATH", os.path.abspath(TEST_DB))
MONTH = "2024-01"


def test_build_deck_returns_valid_pptx():
    """the returned BytesIO should be a re-openable pptx with the expected slide count."""
    buf = build_deck(DB_PATH, month=MONTH, include_ai_commentary=False)
    prs = Presentation(buf)
    assert len(prs.slides) == 6


def test_build_deck_slide_count_same_with_ai_commentary(monkeypatch):
    """AI commentary adds a section to the existing summary slide, not a new slide."""
    monkeypatch.setattr("cli.llm.call_claude", lambda prompt, max_tokens=300: "- insight one\n- insight two")

    buf = build_deck(DB_PATH, month=MONTH, include_ai_commentary=True, api_key="fake-key")
    prs = Presentation(buf)
    assert len(prs.slides) == 6

    summary_slide = prs.slides[1]
    text = "\n".join(
        shape.text_frame.text for shape in summary_slide.shapes if shape.has_text_frame
    )
    assert "Key Insights" in text
    assert "insight one" in text


def test_build_deck_never_calls_llm_without_ai_flag(monkeypatch):
    """build_deck must not hit the network when AI commentary isn't requested."""
    def fail_if_called(*args, **kwargs):
        raise AssertionError("call_claude should not be invoked when include_ai_commentary=False")

    monkeypatch.setattr("cli.llm.call_claude", fail_if_called)
    build_deck(DB_PATH, month=MONTH, include_ai_commentary=False)


def test_kpi_values_match_mart_query():
    """the KPI numbers on the executive summary slide should match a direct mart query."""
    conn = duckdb.connect(DB_PATH, read_only=True)
    row = conn.execute(
        """
        SELECT
            sum(total_trips) as total_trips,
            sum(total_revenue_usd) as total_revenue_usd,
            sum(avg_fare_usd * total_trips) / sum(total_trips) as avg_fare_usd,
            sum(avg_duration_minutes * total_trips) / sum(total_trips) as avg_duration_minutes
        FROM main.mart_daily_kpis
        WHERE strftime(trip_date, '%Y-%m') = ?
        """,
        [MONTH],
    ).fetchone()
    conn.close()
    total_trips, total_revenue_usd, avg_fare_usd, avg_duration_minutes = row

    buf = build_deck(DB_PATH, month=MONTH, include_ai_commentary=False)
    prs = Presentation(buf)
    summary_slide = prs.slides[1]
    text = "\n".join(
        shape.text_frame.text for shape in summary_slide.shapes if shape.has_text_frame
    )

    assert f"{total_trips:,.0f}" in text
    assert f"${total_revenue_usd:,.0f}" in text
    assert f"${avg_fare_usd:,.2f}" in text
    assert f"{avg_duration_minutes:.1f} min" in text


def test_build_deck_raises_for_missing_db():
    with pytest.raises(FileNotFoundError):
        build_deck("/nonexistent/path/to.duckdb", month=MONTH)


def test_build_deck_raises_for_month_with_no_data():
    with pytest.raises(ValueError, match="No data found"):
        build_deck(DB_PATH, month="1999-01")
